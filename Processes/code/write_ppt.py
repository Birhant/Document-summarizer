from pptx import Presentation
from pptx.util import Inches,Pt
from .measure import PPT_page
"""
USE:
create=PPT()
create.write([("BIOGRAPHY","heading"),("This is about birhan","content"),("BIOGRAPHY","heading"),("This is about birhan and how he end up in a great shape after spening ","content"),("a lot of time with anime and manga i think things tgone too far to turn around but their is always aa time.spening a lot of time with anime and manga i think things tgone too far to turn around but their is always aa time.spening a lot of time with anime and manga i think things tgone too far to turn around but their is always aa time.spening a lot of time with anime and manga i think things tgone too far to turn around but their is always aa time.spening a lot of time with anime and manga i think things tgone too far to turn around but their is always aa time.spening a lot of time with anime and manga i think things tgone too far to turn around but their is always aa time.spening a lot of time with anime and manga i think things tgone too far to turn around but their is always aa time and thofidf ethids daca nd tresult. int any akind of trauma this experfj thefiehifn efefhefsb efefbesf eosfahfe fefebfue fuesfbe feufbuebf eience is one of a kind..","content"),])
create.save("Hello")
"""

class PPT:
    def __init__(self):
        self.prs = Presentation()
        self.ppt=PPT_page()

    def title_slide(self,heading,sub_heading):
        subtitle = self.slide.placeholders[1]
        self.title.text = heading
        subtitle.text = sub_heading

    def content_slide(self,heading,sub_heading):
        sub_headings= self.ppt.write(sub_heading)
        self.content_write(heading,sub_headings[0])
        if(len(sub_heading)>1):
            for sub_heading in sub_headings[1:]:
                self.create_slide()
                self.content_write("Cont..",sub_heading)

    def content_write(self,heading,sub_heading):
        subtitle = self.slide.shapes.placeholders[1]
        self.title.text = heading
        subtitle.text=sub_heading
        subtitle.wordwrap=True



    def create_slide(self ,layout=1):
        title_slide_layout = self.prs. slide_layouts[layout]
        self.slide = self.prs. slides. add_slide(title_slide_layout)
        self.title = self.slide.shapes.title

    def write(self,content,cover=True):
        heading=''
        sub_heading=''
        if(cover):
            self.create_slide(0)
            self.title_slide(content[0][0],'summarized text')
            cover=False
        for text,typ in content:
            if(typ.startswith("heading")):
                heading=text
            else:
                sub_heading=text
                self.create_slide(1)
                self.content_slide(heading,sub_heading)

    def save(self,path):
        self.prs.save(path)

